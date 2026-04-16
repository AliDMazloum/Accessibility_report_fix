"""Fix accessibility issues in PowerPoint (.pptx) and Word (.docx) files.

Fixes:
- Set document title and language
- PPTX: ensure all slides have titles
- PPTX: add alt text to images (extract images for external description)

Usage:
    python fix_office.py <input.pptx>              # Fix in place (_fixed.pptx)
    python fix_office.py <input.docx> --scan        # Just show issues
    python fix_office.py <input.pptx> --title "My Slides"
    python fix_office.py <input.pptx> --extract-images /tmp/imgs  # Extract images for alt text
    python fix_office.py <input.pptx> --alt-texts alt.json        # Apply alt texts from JSON
"""
import os
import sys
import json
import argparse
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from docx import Document


def scan_pptx(path):
    """Scan a PPTX for accessibility issues."""
    prs = Presentation(path)
    issues = {}

    if not prs.core_properties.title:
        issues['no_title'] = True
    if not prs.core_properties.language:
        issues['no_language'] = True

    slides_without_title = 0
    images_without_alt = 0
    total_images = 0

    for slide in prs.slides:
        has_title = slide.shapes.title is not None and slide.shapes.title.text.strip()
        if not has_title:
            slides_without_title += 1

        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                total_images += 1
                nv = shape._element.find(
                    './/{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
                if nv is not None:
                    if not nv.get('descr', ''):
                        images_without_alt += 1
                else:
                    images_without_alt += 1

    if slides_without_title > 0:
        issues['slides_without_title'] = slides_without_title
    if images_without_alt > 0:
        issues['images_without_alt'] = images_without_alt
    issues['total_images'] = total_images
    issues['total_slides'] = len(prs.slides)

    return issues


def scan_docx(path):
    """Scan a DOCX for accessibility issues."""
    doc = Document(path)
    issues = {}

    if not doc.core_properties.title:
        issues['no_title'] = True
    if not doc.core_properties.language:
        issues['no_language'] = True

    # Check for headings
    has_headings = False
    for para in doc.paragraphs:
        if para.style.name.startswith('Heading'):
            has_headings = True
            break
    if not has_headings:
        issues['no_headings'] = True

    # Check images
    images_without_alt = 0
    total_images = 0
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            total_images += 1

    issues['total_images'] = total_images

    return issues


def extract_docx_images(path, output_dir):
    """Extract images without alt text from DOCX. Returns list of image info dicts.
    Each dict: {index, image_path, context}"""
    from lxml import etree
    doc = Document(path)
    os.makedirs(output_dir, exist_ok=True)
    images = []

    ns = {
        'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
        'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }

    body = doc.element.body
    drawings = body.findall('.//w:drawing', ns)

    for i, drawing in enumerate(drawings):
        # Find docPr for alt text check
        doc_pr = drawing.find('.//wp:docPr', ns)
        if doc_pr is not None and doc_pr.get('descr', '').strip():
            continue  # Already has alt text

        # Get image data via relationship
        blip = drawing.find('.//a:blip', ns)
        if blip is None:
            continue

        embed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
        if not embed or embed not in doc.part.rels:
            continue

        rel = doc.part.rels[embed]
        img_data = rel.target_part.blob
        ext = rel.target_part.content_type.split('/')[-1]
        if ext == 'jpeg':
            ext = 'jpg'

        img_name = f'image_{i}.{ext}'
        img_path = os.path.join(output_dir, img_name)
        with open(img_path, 'wb') as f:
            f.write(img_data)

        # Get surrounding text for context
        parent_para = drawing.getparent()
        context = ''
        if parent_para is not None:
            # Get text from nearby paragraphs
            all_paras = list(body.findall('.//w:p', ns))
            try:
                idx = all_paras.index(parent_para)
                nearby = all_paras[max(0, idx-2):idx+3]
                for p in nearby:
                    text = ''.join(t.text or '' for t in p.findall('.//w:t', ns))
                    if text.strip():
                        context += text.strip() + ' '
            except ValueError:
                pass

        images.append({
            'index': i,
            'image_path': img_path,
            'context': context.strip()[:500],
        })

    return images


def apply_docx_alt_texts(input_path, output_path, alt_texts):
    """Apply alt texts to images in a DOCX file.
    alt_texts: list of {index, alt_text} dicts."""
    from lxml import etree
    doc = Document(input_path)

    ns = {
        'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
        'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
    }

    # Build lookup
    texts = {}
    for item in alt_texts:
        texts[item['index']] = item['alt_text']

    body = doc.element.body
    drawings = body.findall('.//w:drawing', ns)

    applied = 0
    for i, drawing in enumerate(drawings):
        if i not in texts:
            continue

        doc_pr = drawing.find('.//wp:docPr', ns)
        if doc_pr is not None:
            doc_pr.set('descr', texts[i])
            applied += 1

    if applied > 0:
        doc.save(output_path)

    return {'applied': applied, 'total': len(texts)}


def extract_pptx_images(path, output_dir):
    """Extract images without alt text from PPTX. Returns list of image info dicts.
    Each dict: {slide, shape_idx, image_path, slide_text}"""
    prs = Presentation(path)
    os.makedirs(output_dir, exist_ok=True)
    images = []

    for slide_num, slide in enumerate(prs.slides):
        # Get slide text for context
        slide_text = ''
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text_frame.text + '\n'
        slide_text = slide_text.strip()[:500]

        for shape_idx, shape in enumerate(slide.shapes):
            if not hasattr(shape, 'image'):
                continue

            # Check if it already has alt text
            nv = shape._element.find(
                './/{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
            if nv is not None and nv.get('descr', '').strip():
                continue  # Already has alt text

            # Save image
            image = shape.image
            ext = image.content_type.split('/')[-1]
            if ext == 'jpeg':
                ext = 'jpg'
            img_name = f'slide{slide_num + 1}_shape{shape_idx}.{ext}'
            img_path = os.path.join(output_dir, img_name)

            with open(img_path, 'wb') as f:
                f.write(image.blob)

            images.append({
                'slide': slide_num,
                'shape_idx': shape_idx,
                'image_path': img_path,
                'slide_text': slide_text,
            })

    return images


def apply_pptx_alt_texts(input_path, output_path, alt_texts):
    """Apply alt texts to images in a PPTX file.
    alt_texts: dict mapping 'slideN_shapeM' -> description string,
               or list of {slide, shape_idx, alt_text} dicts."""
    prs = Presentation(input_path)

    # Normalize alt_texts to a dict keyed by (slide, shape_idx)
    texts = {}
    if isinstance(alt_texts, list):
        for item in alt_texts:
            texts[(item['slide'], item['shape_idx'])] = item['alt_text']
    elif isinstance(alt_texts, dict):
        for key, val in alt_texts.items():
            # Parse 'slide1_shape0' format
            parts = key.replace('slide', '').replace('shape', '').split('_')
            if len(parts) == 2:
                texts[(int(parts[0]) - 1, int(parts[1]))] = val

    applied = 0
    for slide_num, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if (slide_num, shape_idx) not in texts:
                continue
            if not hasattr(shape, 'image'):
                continue

            desc = texts[(slide_num, shape_idx)]
            nv = shape._element.find(
                './/{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
            if nv is not None:
                nv.set('descr', desc)
                applied += 1

    if applied > 0:
        prs.save(output_path)

    return {'applied': applied, 'total': len(texts)}


def _pptx_slide_titles_from_semantic(semantic):
    """Return {slide_index: title_text} from semantic.headings entries that
    target whole slides (level==1, locator has slide and no shape_idx, or
    shape_idx is the title placeholder). Otherwise returns empty dict."""
    result = {}
    if not semantic:
        return result
    for h in semantic.get('headings') or []:
        locator = h.get('locator') or {}
        slide = locator.get('slide')
        text = (h.get('text') or '').strip()
        try:
            level = int(h.get('level', 1))
        except Exception:
            level = 1
        if not text or slide is None or level != 1:
            continue
        # First level-1 heading per slide wins.
        if slide not in result:
            result[slide] = text
    return result


def _pptx_table_decision(semantic, slide_idx, shape_idx):
    """Look up a per-table header decision from semantic.tables.
    Returns True / False / None (no decision)."""
    if not semantic:
        return None
    for t in semantic.get('tables') or []:
        loc = t.get('locator') or {}
        if loc.get('slide') == slide_idx and loc.get('shape_idx') == shape_idx:
            v = t.get('has_header_row')
            if isinstance(v, bool):
                return v
    return None


_BULLET_STRIP_RE = None  # lazy compile

def _strip_leading_bullet(text):
    """Remove a leading bullet glyph (• - * or '1.' '2)' '(a)') from a line."""
    import re as _re
    global _BULLET_STRIP_RE
    if _BULLET_STRIP_RE is None:
        _BULLET_STRIP_RE = _re.compile(
            r'^\s*(?:[\u2022\u25E6\u25AA\u25AB\u2219\u00B7\*\-\u2013\u2014]|'
            r'\(?\d+[\.\)]|\(?[a-zA-Z][\.\)])\s+'
        )
    return _BULLET_STRIP_RE.sub('', text, count=1)


def _pptx_apply_list_to_shape(shape, ordered):
    """Apply a real bullet/numbered list to every paragraph in a text frame.

    Writes <a:buAutoNum> (ordered) or <a:buChar> (unordered) to each
    paragraph's pPr and strips leading bullet glyphs from the paragraph text.
    """
    from lxml import etree
    if not shape.has_text_frame:
        return 0
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    tf = shape.text_frame
    applied = 0
    for para in tf.paragraphs:
        p_el = para._pPr if hasattr(para, '_pPr') else None
        pPr = para._p.find(f'{{{ns_a}}}pPr')
        if pPr is None:
            pPr = etree.SubElement(para._p, f'{{{ns_a}}}pPr')
            para._p.insert(0, pPr)
        # Remove existing buChar/buAutoNum/buNone if any.
        for tag in ('buChar', 'buAutoNum', 'buNone', 'buBlip'):
            for el in pPr.findall(f'{{{ns_a}}}{tag}'):
                pPr.remove(el)
        if ordered:
            bu = etree.SubElement(pPr, f'{{{ns_a}}}buAutoNum')
            bu.set('type', 'arabicPeriod')
        else:
            bu = etree.SubElement(pPr, f'{{{ns_a}}}buChar')
            bu.set('char', '\u2022')
        # Strip leading glyph from the paragraph's runs' text.
        for run in para.runs:
            if run.text:
                new_text = _strip_leading_bullet(run.text)
                if new_text != run.text:
                    run.text = new_text
                break  # only first run of paragraph carries the glyph
        applied += 1
    return applied


def fix_pptx(input_path, output_path=None, title=None, lang='en-US',
             alt_texts=None, semantic=None):
    """Fix accessibility issues in a PPTX file.

    When ``semantic`` is provided (dict from semantic_agents), its values are
    preferred over deterministic fallbacks:
      - semantic.semantic_title / language set core properties unconditionally
        when present (overwrites existing values).
      - semantic.headings with level==1 + {slide: N} provide slide titles.
      - semantic.tables { has_header_row: bool } gate the firstRow flag per
        table (FALSE means: do not mark row 0 as header).
      - semantic.lists convert manually-bulleted paragraphs into real
        bullet/numbered lists on the identified shape.
    """
    if output_path is None:
        stem, ext = os.path.splitext(input_path)
        if stem.endswith('_fixed'):
            output_path = input_path
        else:
            output_path = stem + '_fixed' + ext

    prs = Presentation(input_path)
    fixed = []

    # Title: semantic wins if provided; otherwise explicit arg; else filename.
    sem_title = (semantic or {}).get('semantic_title')
    sem_lang = (semantic or {}).get('language')

    if sem_title:
        if prs.core_properties.title != sem_title:
            prs.core_properties.title = sem_title
            fixed.append('title (semantic)')
    else:
        doc_title = title or os.path.splitext(os.path.basename(input_path))[0].replace('_', ' ')
        if not prs.core_properties.title:
            prs.core_properties.title = doc_title
            fixed.append('title')

    # Language.
    if sem_lang:
        if prs.core_properties.language != sem_lang:
            prs.core_properties.language = sem_lang
            fixed.append('language (semantic)')
    elif not prs.core_properties.language:
        prs.core_properties.language = lang
        fixed.append('language')

    # Ensure all slides have titles. Prefer semantic-provided titles when
    # available; fall back to largest-text heuristic.
    semantic_slide_titles = _pptx_slide_titles_from_semantic(semantic)
    slides_fixed = 0
    semantic_titles_used = 0
    for i, slide in enumerate(prs.slides):
        has_title = slide.shapes.title is not None and slide.shapes.title.text.strip()
        if has_title:
            continue

        chosen_text = semantic_slide_titles.get(i)
        if not chosen_text:
            # Largest-text heuristic (original fallback).
            best_text = ''
            best_size = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.font.size > best_size:
                                best_size = run.font.size
                                best_text = para.text.strip()
                            elif not run.font.size and para.text.strip() and not best_text:
                                best_text = para.text.strip()
            chosen_text = best_text
        else:
            semantic_titles_used += 1

        if slide.shapes.title is not None and chosen_text:
            slide.shapes.title.text = chosen_text
            slides_fixed += 1

    if slides_fixed > 0:
        label = f'{slides_fixed} slide titles'
        if semantic_titles_used > 0:
            label += f' ({semantic_titles_used} semantic)'
        fixed.append(label)

    # Table headers: respect semantic decision if present, else default to TRUE
    # (preserves prior behavior for documents without semantic validation).
    tables_fixed = 0
    tables_left_alone = 0
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_table:
                continue
            decision = _pptx_table_decision(semantic, slide_idx, shape_idx)
            if decision is False:
                tables_left_alone += 1
                continue  # Validator says: not a header row. Don't mark it.
            # decision is True or None — proceed to set firstRow=1.
            tbl_el = shape.table._tbl
            ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            tblPr = tbl_el.find(f'{{{ns_a}}}tblPr')
            if tblPr is None:
                from lxml import etree
                tblPr = etree.SubElement(tbl_el, f'{{{ns_a}}}tblPr')
                tbl_el.insert(0, tblPr)
            if tblPr.get('firstRow') != '1':
                tblPr.set('firstRow', '1')
                tblPr.set('bandRow', '1')
                tables_fixed += 1

    if tables_fixed > 0:
        fixed.append(f'{tables_fixed} table headers')
    if tables_left_alone > 0:
        fixed.append(f'{tables_left_alone} tables left unflagged (semantic)')

    # Lists: apply semantic list decisions.
    if semantic and semantic.get('lists'):
        lists_applied = 0
        for ls in semantic['lists']:
            loc = ls.get('locator') or {}
            slide_idx = loc.get('slide')
            shape_idx = loc.get('shape_idx')
            if slide_idx is None or shape_idx is None:
                continue
            if slide_idx < 0 or slide_idx >= len(prs.slides):
                continue
            slide = prs.slides[slide_idx]
            shapes = list(slide.shapes)
            if shape_idx < 0 or shape_idx >= len(shapes):
                continue
            try:
                applied = _pptx_apply_list_to_shape(shapes[shape_idx], bool(ls.get('ordered')))
                if applied > 0:
                    lists_applied += 1
            except Exception:
                continue
        if lists_applied > 0:
            fixed.append(f'{lists_applied} lists (semantic)')

    # Save intermediate if we have changes so far, or just copy for alt text step
    if fixed or alt_texts:
        prs.save(output_path)

    # Apply alt texts if provided (must be done after save since it re-opens the file)
    if alt_texts:
        result = apply_pptx_alt_texts(output_path if os.path.exists(output_path) else input_path,
                                      output_path, alt_texts)
        if result['applied'] > 0:
            fixed.append(f"{result['applied']} image alt texts")

    if not fixed:
        return {'status': 'no issues found'}

    return {
        'status': 'fixed',
        'fixed': fixed,
        'output': output_path,
        'size': os.path.getsize(output_path),
    }


def _docx_neutralize_heading_style(style):
    """Strip visual run-property defaults from a Word Heading style so applying
    it to a paragraph does not change its visible font, color, or size.

    Word's built-in "Heading 1/2/3" styles define a themed font (Calibri Light),
    a teal color, and a larger size. When we apply the style to a normal
    paragraph for a11y tagging, we don't want those visuals to override the
    body text. We remove the <w:rPr> run-properties block on the style.
    """
    try:
        from docx.oxml.ns import qn
        from lxml import etree
        el = style.element
        rPr = el.find(qn('w:rPr'))
        if rPr is not None:
            el.remove(rPr)
    except Exception:
        pass


def _docx_get_or_create_heading_style(doc, level):
    name = f'Heading {level}'
    try:
        style = doc.styles[name]
    except KeyError:
        style = doc.styles.add_style(name, 1)  # 1 = paragraph
        style.base_style = doc.styles['Normal']
    # Strip visuals so the style is an invisible semantic tag.
    _docx_neutralize_heading_style(style)
    return style


def _docx_get_or_create_list_style(doc, ordered):
    """Return the 'List Number' or 'List Bullet' style, creating if missing."""
    name = 'List Number' if ordered else 'List Bullet'
    try:
        return doc.styles[name]
    except KeyError:
        style = doc.styles.add_style(name, 1)
        style.base_style = doc.styles['Normal']
        return style


def _docx_resolve_style_font(style):
    """Walk a style's base_style chain and return the first non-None value
    for each font property. python-docx resolves style inheritance lazily;
    we need the effective values before we flip the paragraph to a new style."""
    name = size = bold = italic = underline = color = None
    cur = style
    seen = set()
    while cur is not None and id(cur) not in seen:
        seen.add(id(cur))
        try:
            f = cur.font
        except Exception:
            break
        if name is None and f.name:
            name = f.name
        if size is None and f.size:
            size = f.size
        if bold is None and f.bold is not None:
            bold = f.bold
        if italic is None and f.italic is not None:
            italic = f.italic
        if underline is None and f.underline is not None:
            underline = f.underline
        if color is None:
            try:
                if f.color and f.color.rgb is not None:
                    color = f.color.rgb
            except Exception:
                pass
        cur = getattr(cur, 'base_style', None)
    return {
        'name': name, 'size': size, 'bold': bold,
        'italic': italic, 'underline': underline, 'color': color,
    }


def _docx_snapshot_runs(para):
    """Capture the effective run formatting (direct + inherited from the
    paragraph's current style) so we can restore the visible look after
    flipping the paragraph to a new style like Heading N or List Bullet."""
    style_defaults = _docx_resolve_style_font(para.style) if para.style else {}
    snapshot = []
    for run in para.runs:
        f = run.font
        color_rgb = None
        try:
            if f.color and f.color.rgb is not None:
                color_rgb = f.color.rgb
        except Exception:
            color_rgb = None
        # For each property, prefer direct run value; fall back to the style
        # chain's effective value so we can re-impose it after the style flip.
        snapshot.append({
            'name': f.name if f.name else style_defaults.get('name'),
            'size': f.size if f.size else style_defaults.get('size'),
            'bold': f.bold if f.bold is not None else style_defaults.get('bold'),
            'italic': f.italic if f.italic is not None else style_defaults.get('italic'),
            'underline': f.underline if f.underline is not None else style_defaults.get('underline'),
            'color': color_rgb if color_rgb is not None else style_defaults.get('color'),
        })
    return snapshot


def _docx_restore_runs(para, snapshot):
    """Reapply captured direct run formatting after a style change."""
    for run, snap in zip(para.runs, snapshot):
        f = run.font
        if snap['name']:
            f.name = snap['name']
        if snap['size']:
            f.size = snap['size']
        if snap['bold'] is not None:
            f.bold = snap['bold']
        if snap['italic'] is not None:
            f.italic = snap['italic']
        if snap['underline'] is not None:
            f.underline = snap['underline']
        if snap['color'] is not None:
            try:
                f.color.rgb = snap['color']
            except Exception:
                pass


def _docx_apply_semantic_headings(doc, semantic):
    """Apply Heading N style to paragraphs at given locator.para_index.
    Preserves the paragraph's original run formatting so the visual look
    does not flip to Word's default heading font/color. Returns count applied.
    """
    paragraphs = doc.paragraphs
    count = 0
    for h in semantic.get('headings') or []:
        locator = h.get('locator') or {}
        idx = locator.get('para_index')
        if not isinstance(idx, int) or idx < 0 or idx >= len(paragraphs):
            continue
        try:
            level = int(h.get('level', 1))
        except Exception:
            level = 1
        level = max(1, min(3, level))
        para = paragraphs[idx]
        snapshot = _docx_snapshot_runs(para)
        para.style = _docx_get_or_create_heading_style(doc, level)
        _docx_restore_runs(para, snapshot)
        count += 1
    return count


def _docx_apply_semantic_lists(doc, semantic):
    """Apply List Bullet / List Number style across the paragraph ranges in
    semantic.lists; strip leading bullet glyphs from the first run."""
    paragraphs = doc.paragraphs
    applied = 0
    for ls in semantic.get('lists') or []:
        locator = ls.get('locator') or {}
        para_range = locator.get('para_range')
        if not (isinstance(para_range, (list, tuple)) and len(para_range) == 2):
            continue
        lo, hi = para_range
        if not (isinstance(lo, int) and isinstance(hi, int)) or lo > hi:
            continue
        lo = max(0, lo)
        hi = min(len(paragraphs) - 1, hi)
        ordered = bool(ls.get('ordered'))
        style = _docx_get_or_create_list_style(doc, ordered)
        for pi in range(lo, hi + 1):
            para = paragraphs[pi]
            snapshot = _docx_snapshot_runs(para)
            para.style = style
            _docx_restore_runs(para, snapshot)
            # Strip leading glyph from the first run's text.
            for run in para.runs:
                if run.text:
                    new_text = _strip_leading_bullet(run.text)
                    if new_text != run.text:
                        run.text = new_text
                    break
            applied += 1
    return applied


def _docx_table_decision(semantic, table_idx):
    if not semantic:
        return None
    for t in semantic.get('tables') or []:
        if t.get('table_idx') == table_idx:
            v = t.get('has_header_row')
            if isinstance(v, bool):
                return v
    return None


def fix_docx(input_path, output_path=None, title=None, lang='en-US', semantic=None):
    """Fix accessibility issues in a DOCX file.

    When ``semantic`` is provided:
      - semantic.semantic_title / language overwrite core properties.
      - semantic.headings replace the font-size/pattern heading detection.
      - semantic.tables { has_header_row } gate the w:tblHeader flag per table.
      - semantic.lists turn manually-bulleted paragraph ranges into real
        List Bullet / List Number paragraphs.
    """
    if output_path is None:
        stem, ext = os.path.splitext(input_path)
        if stem.endswith('_fixed'):
            output_path = input_path
        else:
            output_path = stem + '_fixed' + ext

    doc = Document(input_path)
    fixed = []

    sem_title = (semantic or {}).get('semantic_title')
    sem_lang = (semantic or {}).get('language')

    # Title
    if sem_title:
        if doc.core_properties.title != sem_title:
            doc.core_properties.title = sem_title
            fixed.append('title (semantic)')
    else:
        doc_title = title or os.path.splitext(os.path.basename(input_path))[0].replace('_', ' ')
        if not doc.core_properties.title:
            doc.core_properties.title = doc_title
            fixed.append('title')

    # Language
    if sem_lang:
        if doc.core_properties.language != sem_lang:
            doc.core_properties.language = sem_lang
            fixed.append('language (semantic)')
    elif not doc.core_properties.language:
        doc.core_properties.language = lang
        fixed.append('language')

    # Headings: prefer semantic list; otherwise fall back to heuristic path.
    use_semantic_headings = bool(semantic and (semantic.get('headings') or []))
    if use_semantic_headings:
        n = _docx_apply_semantic_headings(doc, semantic)
        if n > 0:
            fixed.append(f'{n} headings (semantic)')

    # Heuristic heading detection (retained as fallback when no semantic input
    # and document currently has no headings).
    has_headings = any(p.style.name.startswith('Heading') for p in doc.paragraphs)
    if not use_semantic_headings and not has_headings:
        headings_added = 0

        # Method 1: Detect by font size hierarchy
        font_sizes = {}
        for para in doc.paragraphs:
            for run in para.runs:
                if run.font.size and para.text.strip():
                    size = run.font.size
                    if size not in font_sizes:
                        font_sizes[size] = 0
                    font_sizes[size] += 1

        if len(font_sizes) > 1:
            body_size = max(font_sizes, key=font_sizes.get)
            heading_sizes = sorted([s for s in font_sizes if s > body_size], reverse=True)

            heading_map = {}
            for i, size in enumerate(heading_sizes[:3]):
                heading_map[size] = f'Heading {i + 1}'

            # Ensure heading styles exist
            def get_heading_style(doc, level):
                name = f'Heading {level}'
                try:
                    return doc.styles[name]
                except KeyError:
                    from docx.shared import Pt as DPt
                    style = doc.styles.add_style(name, 1)  # 1 = paragraph
                    style.base_style = doc.styles['Normal']
                    style.font.bold = True
                    style.font.size = DPt(16 - (level - 1) * 2)
                    return style

            for para in doc.paragraphs:
                if not para.text.strip():
                    continue
                max_size = 0
                for run in para.runs:
                    if run.font.size and run.font.size > max_size:
                        max_size = run.font.size
                if max_size in heading_map:
                    level = int(heading_map[max_size].split()[-1])
                    snapshot = _docx_snapshot_runs(para)
                    para.style = get_heading_style(doc, level)
                    _docx_restore_runs(para, snapshot)
                    headings_added += 1

        # Method 2: If no font size hierarchy, detect by patterns
        # (short bold text, text ending with ":", numbered sections)
        if headings_added == 0:
            import re

            def get_heading_style_safe(doc, level=1):
                name = f'Heading {level}'
                try:
                    return doc.styles[name]
                except KeyError:
                    from docx.shared import Pt as DPt
                    style = doc.styles.add_style(name, 1)
                    style.base_style = doc.styles['Normal']
                    style.font.bold = True
                    style.font.size = DPt(16 - (level - 1) * 2)
                    return style

            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                is_bold = all(r.bold for r in para.runs if r.text.strip()) and para.runs
                is_short = len(text) < 80
                looks_like_heading = (
                    (is_bold and is_short) or
                    re.match(r'^(Exercise|Question|Problem|Part|Section|Chapter)\s+\d', text, re.I) or
                    (is_short and text.endswith(':') and len(text) < 40)
                )
                if looks_like_heading:
                    snapshot = _docx_snapshot_runs(para)
                    para.style = get_heading_style_safe(doc)
                    _docx_restore_runs(para, snapshot)
                    headings_added += 1

        if headings_added > 0:
            fixed.append(f'{headings_added} headings')

    # Fix table headers: mark first row as w:tblHeader unless the validator
    # says this table doesn't have a header row. Only body-level tables are
    # indexed (matches docx.tables ordering); nested tables still default to
    # TRUE to preserve prior behavior.
    from lxml import etree
    from docx.oxml.ns import qn
    tables_fixed = 0
    tables_left_alone = 0

    body_tables = list(doc.tables)
    body_tbl_elements = {id(t._tbl) for t in body_tables}

    all_tables = doc.element.body.findall('.//' + qn('w:tbl'))
    body_idx = 0
    for tbl in all_tables:
        is_body_tbl = id(tbl) in body_tbl_elements
        decision = None
        if is_body_tbl:
            decision = _docx_table_decision(semantic, body_idx)
            body_idx += 1
        if decision is False:
            tables_left_alone += 1
            continue
        rows = tbl.findall(qn('w:tr'))
        if not rows:
            continue
        first_row = rows[0]
        trPr = first_row.find(qn('w:trPr'))
        if trPr is None:
            trPr = etree.SubElement(first_row, qn('w:trPr'))
            first_row.insert(0, trPr)
        tblHeader = trPr.find(qn('w:tblHeader'))
        if tblHeader is None:
            etree.SubElement(trPr, qn('w:tblHeader'))
            tables_fixed += 1

    if tables_fixed > 0:
        fixed.append(f'{tables_fixed} table headers')
    if tables_left_alone > 0:
        fixed.append(f'{tables_left_alone} tables left unflagged (semantic)')

    # Lists: apply semantic list styling where ranges are provided.
    if semantic and (semantic.get('lists') or []):
        lists_applied = _docx_apply_semantic_lists(doc, semantic)
        if lists_applied > 0:
            fixed.append(f'{lists_applied} list items (semantic)')

    if not fixed:
        return {'status': 'no issues found'}

    doc.save(output_path)

    return {
        'status': 'fixed',
        'fixed': fixed,
        'output': output_path,
        'size': os.path.getsize(output_path),
    }


def fix_office(input_path, output_path=None, title=None, lang='en-US', semantic=None):
    """Fix any Office file based on extension."""
    ext = os.path.splitext(input_path)[1].lower()
    if ext == '.pptx':
        return fix_pptx(input_path, output_path, title, lang, semantic=semantic)
    elif ext == '.docx':
        return fix_docx(input_path, output_path, title, lang, semantic=semantic)
    else:
        return {'status': f'unsupported format: {ext}'}


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Fix Office file accessibility')
    parser.add_argument('input', help='PPTX or DOCX file')
    parser.add_argument('-o', '--output', help='Output path')
    parser.add_argument('--title', help='Document title')
    parser.add_argument('--lang', default='en-US', help='Language')
    parser.add_argument('--scan', action='store_true', help='Just scan for issues')
    parser.add_argument('--extract-images', metavar='DIR',
                        help='Extract images without alt text to DIR (PPTX only)')
    parser.add_argument('--alt-texts', metavar='JSON',
                        help='Apply alt texts from JSON file (PPTX only)')
    args = parser.parse_args()

    ext = os.path.splitext(args.input)[1].lower()

    if args.extract_images:
        if ext != '.pptx':
            print('--extract-images only works with PPTX files')
            sys.exit(1)
        images = extract_pptx_images(args.input, args.extract_images)
        print(f'Extracted {len(images)} images without alt text:')
        for img in images:
            print(f"  Slide {img['slide']+1}, shape {img['shape_idx']}: {img['image_path']}")
        # Save metadata
        meta_path = os.path.join(args.extract_images, 'images_meta.json')
        with open(meta_path, 'w') as f:
            json.dump(images, f, indent=2)
        print(f'Metadata saved to {meta_path}')
    elif args.scan:
        if ext == '.pptx':
            issues = scan_pptx(args.input)
        elif ext == '.docx':
            issues = scan_docx(args.input)
        else:
            print(f'Unsupported: {ext}')
            sys.exit(1)

        print(f'Issues found:')
        for k, v in issues.items():
            print(f'  {k}: {v}')
    else:
        alt_texts = None
        if args.alt_texts:
            with open(args.alt_texts) as f:
                alt_texts = json.load(f)

        if ext == '.pptx':
            result = fix_pptx(args.input, args.output, args.title, args.lang, alt_texts)
        elif ext == '.docx':
            result = fix_docx(args.input, args.output, args.title, args.lang)
        else:
            result = fix_office(args.input, args.output, args.title, args.lang)

        if result['status'] == 'fixed':
            print(f"Fixed: {', '.join(result['fixed'])}")
            print(f"Output: {result['output']} ({result['size']:,} bytes)")
        else:
            print(f"Status: {result['status']}")
